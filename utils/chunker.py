import os
import json
import email
from email import policy
import requests
from bs4 import BeautifulSoup
from typing import List, Dict, Any

from agno.knowledge.document import Document
from agno.knowledge.chunking.recursive import RecursiveChunking

def extract_text_from_file(file_path: str) -> str:
    """Extracts raw text from various enterprise file formats."""
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        import pypdf
        text = ""
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted: text += extracted + "\n"
        return text
    
    elif ext == ".docx":
        import docx
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
        
    elif ext == ".pptx":
        import pptx
        prs = pptx.Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_runs.append(shape.text)
        return "\n\n".join(text_runs)
        
    elif ext == ".ppt":
         raise ValueError("Legacy .ppt binary files are not natively supported. Please convert to .pptx format.")

    elif ext == ".xml":
        with open(file_path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "xml")
            return soup.get_text(separator=" ")

    elif ext == ".eml":
        with open(file_path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=policy.default)
        content = f"Subject: {msg.get('subject', '')}\nFrom: {msg.get('from', '')}\nTo: {msg.get('to', '')}\n\n"
        body = msg.get_body(preferencelist=('plain',))
        if body:
            content += body.get_content()
        return content

    elif ext in [".txt", ".md"]:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
            
    else:
        raise ValueError(f"Unsupported file extension for plain text extraction: {ext}")

def extract_text_from_url(url: str) -> str:
    """Extracts clean text from a web page."""
    response = requests.get(url, timeout=10)
    response.raise_for_status()
    soup = BeautifulSoup(response.text, "html.parser")
    
    for script in soup(["script", "style"]):
        script.extract()
        
    text = soup.get_text(separator=" ")
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    return '\n'.join(chunk for chunk in chunks if chunk)

def chunk_text(
    text: str,
    chunk_size: int = 1000,
    overlap: int = 200,
    source_metadata: str = "text"
) -> List[Dict[str, Any]]:
    """Chunks text uniformly using Agno's abstract RecursiveChunking."""
    if not text or not text.strip():
        return []

    splitter = RecursiveChunking(
        chunk_size=chunk_size,
        overlap=overlap,
    )
    
    doc = Document(content=text, name=source_metadata)
    chunked_docs = splitter.chunk(doc)
    
    return [
        {
            "chunk_index": i,
            "text": d.content,
            "source": source_metadata
        }
        for i, d in enumerate(chunked_docs)
    ]

def chunk_document(file_path: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
    ext = os.path.splitext(file_path)[1].lower()
    source_metadata = os.path.basename(file_path)

    # Context-aware tabular chunking for spreadsheets / CSVs
    if ext in [".xlsx", ".xls", ".csv"]:
        import pandas as pd
        if ext == ".csv":
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
            
        chunks = []
        header = df.columns.tolist()
        current_chunk_text = ""
        
        for _, row in df.iterrows():
            row_str = " | ".join(f"{h}: {row[h]}" for h in header if pd.notna(row[h])) + "\n"
            if len(current_chunk_text) + len(row_str) > chunk_size and current_chunk_text:
                chunks.append({
                    "chunk_index": len(chunks),
                    "text": current_chunk_text.strip(),
                    "source": source_metadata
                })
                current_chunk_text = row_str 
            else:
                current_chunk_text += row_str
                
        if current_chunk_text:
            chunks.append({
                "chunk_index": len(chunks),
                "text": current_chunk_text.strip(),
                "source": source_metadata
            })
        return chunks

    # Context-aware chunking for JSON APIs / configurations
    elif ext == ".json":
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
            
        # If it's an array of distinct objects, chunk object-by-object seamlessly
        if isinstance(data, list) and all(isinstance(item, dict) for item in data):
            chunks = []
            current_chunk_text = ""
            for item in data:
                item_str = json.dumps(item) + "\n"
                if len(current_chunk_text) + len(item_str) > chunk_size and current_chunk_text:
                    chunks.append({
                        "chunk_index": len(chunks),
                        "text": current_chunk_text.strip(),
                        "source": source_metadata
                    })
                    current_chunk_text = item_str
                else:
                    current_chunk_text += item_str
                    
            if current_chunk_text:
                chunks.append({
                    "chunk_index": len(chunks),
                    "text": current_chunk_text.strip(),
                    "source": source_metadata
                })
            return chunks
        else:
            # Fallback to recursively chunking the whole JSON blob
            text = json.dumps(data, indent=2)
            return chunk_text(text, chunk_size, overlap, source_metadata=source_metadata)

    # Standard recursive chunking for regular documents
    text = extract_text_from_file(file_path)
    return chunk_text(text, chunk_size, overlap, source_metadata=source_metadata)

def chunk_webpage(url: str, chunk_size: int = 1000, overlap: int = 200) -> List[Dict[str, Any]]:
    text = extract_text_from_url(url)
    return chunk_text(text, chunk_size, overlap, source_metadata=url)
